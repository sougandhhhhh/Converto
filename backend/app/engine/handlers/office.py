import os
import shutil
import zipfile
import logging
import tempfile
from pathlib import Path
import pandas as pd
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from app.engine.handlers.base import BaseHandler

logger = logging.getLogger(__name__)

class OfficeHandler(BaseHandler):
    """
    Handles all Office-related conversions (DOCX, PPTX, XLSX).
    Utilizes LibreOffice headless for high-fidelity conversions,
    and Pandas/openpyxl for spreadsheet mutations (CSV, JSON, XML).
    """

    def convert(self, input_path: Path, output_path: Path, from_ext: str, to_ext: str) -> Path:
        from_ext = from_ext.lower()
        to_ext = to_ext.lower()
        output_dir = output_path.parent

        logger.info(f"OfficeHandler converting {from_ext} -> {to_ext}")

        # 1. Image Targets (JPG, PNG, WEBP, HEIC)
        # Convert Office to PDF first, then use PDFHandler to extract images
        if to_ext in [".jpg", ".jpeg", ".png", ".webp", ".heic"]:
            temp_pdf = input_path.with_suffix(".pdf")
            try:
                self._convert_with_libreoffice(input_path, temp_pdf, "pdf")
                from app.engine.handlers.pdf import PDFHandler
                pdf_h = PDFHandler()
                return pdf_h.convert(temp_pdf, output_path, ".pdf", to_ext)
            finally:
                if temp_pdf.exists():
                    try:
                        temp_pdf.unlink()
                    except Exception as e:
                        logger.warning(f"Failed to delete temp PDF: {e}")

        # 2. ZIP target
        if to_ext == ".zip":
            with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zip_file:
                zip_file.write(input_path, arcname=input_path.name)
            return output_path

        # 3. Spreadsheet mutations (XLSX -> CSV, JSON, XML, TXT, HTML)
        if from_ext == ".xlsx":
            if to_ext == ".csv":
                df = pd.read_excel(input_path)
                df.to_csv(output_path, index=False)
                return output_path
            elif to_ext == ".json":
                df = pd.read_excel(input_path)
                df.to_json(output_path, orient="records", indent=4)
                return output_path
            elif to_ext == ".xml":
                df = pd.read_excel(input_path)
                df.columns = [f"col_{idx}" for idx, _ in enumerate(df.columns, start=1)]
                df.to_xml(output_path, index=False)
                return output_path
            elif to_ext == ".txt":
                # Convert spreadsheet to tab-delimited text
                df = pd.read_excel(input_path)
                df.to_csv(output_path, sep="\t", index=False)
                return output_path
            elif to_ext == ".html":
                df = pd.read_excel(input_path)
                df.to_html(output_path, index=False, classes="table table-striped")
                return output_path

        # 4. Semantic cross-format fallbacks for routes LibreOffice doesn't reliably support
        if from_ext == ".docx" and to_ext == ".xlsx":
            return self._docx_to_xlsx(input_path, output_path)
        if from_ext == ".docx" and to_ext == ".pptx":
            return self._docx_to_pptx(input_path, output_path)
        if from_ext == ".pptx" and to_ext == ".docx":
            return self._pptx_to_docx(input_path, output_path)
        if from_ext == ".pptx" and to_ext == ".txt":
            return self._pptx_to_txt(input_path, output_path)
        if from_ext == ".pptx" and to_ext == ".html":
            return self._pptx_to_html(input_path, output_path)
        if from_ext == ".xlsx" and to_ext == ".docx":
            return self._xlsx_to_docx(input_path, output_path)

        # 5. Standard LibreOffice conversions
        # Map target extension to LibreOffice filter formats
        format_map = {
            ".pdf": "pdf:writer_pdf_Export",
            ".docx": "docx:MS Word 2007 XML",
            ".pptx": "pptx:Impress MS PowerPoint 2007 XML",
            ".xlsx": "xlsx:Calc MS Excel 2007 XML",
            ".txt": "txt:Text",
            ".csv": "csv:Text - txt - csv (StarCalc)",
            ".html": "html:XHTML Writer File",
        }

        lo_format = format_map.get(to_ext)
        if not lo_format:
            raise ValueError(f"No LibreOffice format mapping for target: {to_ext}")

        self._convert_with_libreoffice(input_path, output_path, lo_format)
        return output_path

    def _convert_with_libreoffice(self, input_path: Path, output_path: Path, lo_format: str):
        """
        Executes a LibreOffice headless conversion command.
        """
        temp_outdir = output_path.parent
        with tempfile.TemporaryDirectory(prefix="lo-profile-") as profile_dir:
            profile_uri = Path(profile_dir).resolve().as_uri()
            cmd = [
                "libreoffice",
                "--headless",
                f"-env:UserInstallation={profile_uri}",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--norestore",
                "--convert-to",
                lo_format,
                "--outdir",
                str(temp_outdir),
                str(input_path),
            ]

            logger.info(f"Executing LibreOffice: {' '.join(cmd)}")
            self.run_subprocess(cmd, timeout=120)

        # LibreOffice names the output file by replacing the input suffix with the target suffix
        expected_suffix = output_path.suffix
        expected_filename = f"{input_path.stem}{expected_suffix}"
        lo_output_file = temp_outdir / expected_filename

        if lo_output_file.exists():
            if lo_output_file != output_path:
                shutil.move(str(lo_output_file), str(output_path))
            logger.info(f"LibreOffice conversion successful: {output_path}")
        else:
            # Look for any case differences or double suffixes
            possible_files = list(temp_outdir.glob(f"{input_path.stem}.*"))
            # Filter by matching target suffix
            matching_files = [f for f in possible_files if f.suffix.lower() == expected_suffix.lower()]
            if matching_files:
                shutil.move(str(matching_files[0]), str(output_path))
                logger.info(f"LibreOffice conversion successful (resolved suffix match): {output_path}")
            else:
                raise FileNotFoundError(
                    f"LibreOffice execution completed but target file was not found at {lo_output_file}"
                )

    def _docx_to_xlsx(self, input_path: Path, output_path: Path) -> Path:
        doc = Document(input_path)
        wb = Workbook()
        ws = wb.active
        ws.title = "Document"
        ws.cell(row=1, column=1, value="Content")
        row_idx = 2
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                ws.cell(row=row_idx, column=1, value=text)
                row_idx += 1
        wb.save(output_path)
        return output_path

    def _docx_to_pptx(self, input_path: Path, output_path: Path) -> Path:
        doc = Document(input_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        prs = Presentation()
        chunk_size = 8
        for start in range(0, max(len(paragraphs), 1), chunk_size):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = input_path.stem
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            chunk = paragraphs[start : start + chunk_size] or ["(No extractable text)"]
            for idx, line in enumerate(chunk):
                if idx == 0:
                    body.text = line
                else:
                    p = body.add_paragraph()
                    p.text = line
        prs.save(output_path)
        return output_path

    def _pptx_to_docx(self, input_path: Path, output_path: Path) -> Path:
        prs = Presentation(input_path)
        doc = Document()
        for slide_index, slide in enumerate(prs.slides, start=1):
            heading = doc.add_paragraph()
            heading.add_run(f"Slide {slide_index}").bold = True
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        doc.add_paragraph(text)
        doc.save(output_path)
        return output_path

    def _pptx_to_txt(self, input_path: Path, output_path: Path) -> Path:
        prs = Presentation(input_path)
        lines = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            lines.append(f"Slide {slide_index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        lines.append(text)
            lines.append("")
        output_path.write_text("\n".join(lines).strip(), encoding="utf-8")
        return output_path

    def _xlsx_to_docx(self, input_path: Path, output_path: Path) -> Path:
        df = pd.read_excel(input_path, dtype=str).fillna("")
        doc = Document()
        if df.empty:
            doc.add_paragraph("(No spreadsheet content)")
        else:
            doc.add_paragraph("Spreadsheet Export").runs[0].bold = True
            doc.add_paragraph("")
            header = " | ".join([str(col) for col in df.columns])
            doc.add_paragraph(header).runs[0].bold = True
            for _, row in df.iterrows():
                doc.add_paragraph(" | ".join([str(v) for v in row.tolist()]))
        doc.save(output_path)
        return output_path

    def _pptx_to_html(self, input_path: Path, output_path: Path) -> Path:
        prs = Presentation(input_path)
        sections = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        lines.append(text)
            body = "".join(f"<p>{line}</p>" for line in lines) if lines else "<p>(No text content)</p>"
            sections.append(f"<section><h2>Slide {slide_index}</h2>{body}</section>")
        html = (
            "<!DOCTYPE html><html><head><meta charset='utf-8'>"
            f"<title>{input_path.stem}</title>"
            "<style>body{font-family:Arial,sans-serif;margin:24px;}section{margin-bottom:18px;}h2{margin-bottom:8px;}</style>"
            "</head><body>"
            + "".join(sections)
            + "</body></html>"
        )
        output_path.write_text(html, encoding="utf-8")
        return output_path
