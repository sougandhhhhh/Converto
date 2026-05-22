import os
import shutil
import zipfile
import logging
import re
import html
import base64
import mimetypes
from datetime import datetime
from pathlib import Path
from PIL import Image
import pdfplumber
import camelot
from docx import Document
from docx.shared import Inches as DocxInches
from pptx import Presentation
from pptx.util import Inches, Pt
from app.engine.handlers.base import BaseHandler
try:
    import pillow_heif
    pillow_heif.register_heif_opener()
    HEIF_PLUGIN_ENABLED = True
except Exception:
    HEIF_PLUGIN_ENABLED = False

logger = logging.getLogger(__name__)
PDF_IMAGE_RENDER_TIMEOUT_SECONDS = 300
QUALITY_PRESETS = {
    "fast": {
        "image_dpi": 135,
        "pptx_dpi": 140,
        "docx_fallback_dpi": 140,
        "zip_dpi": 100,
    },
    "high": {
        "image_dpi": 300,
        "pptx_dpi": 300,
        "docx_fallback_dpi": 300,
        "zip_dpi": 220,
    },
}

class PDFHandler(BaseHandler):
    """
    Handles high-fidelity PDF conversions to documents, presentations,
    spreadsheets, images, text, HTML, HEIC, and packaged ZIPs.
    Integrates OCRmyPDF for scanned inputs, Camelot for table mapping,
    and python-pptx / python-docx for semantic reconstructions.
    """

    def convert(self, input_path: Path, output_path: Path, from_ext: str, to_ext: str, quality: str = "fast") -> Path:
        to_ext = to_ext.lower()
        image_targets = {".jpg", ".jpeg", ".png", ".webp", ".heic"}
        quality = self._normalize_quality(quality)

        logger.info(f"PDFHandler converting PDF -> {to_ext} quality={quality}")
        # Fast path for image outputs: avoid expensive normalization/OCR heuristics.
        if to_ext in image_targets:
            return self._to_image(input_path, output_path, to_ext, quality=quality)

        pdf_source = self._normalize_pdf(input_path)

        # Detect if PDF is scanned (little to no extractable text)
        is_scanned = self._is_scanned_pdf(pdf_source)

        if is_scanned:
            logger.info("Scanned PDF detected. Triggering OCRmyPDF pipeline...")
            ocr_pdf = input_path.parent / f"{input_path.stem}_ocr.pdf"
            try:
                # OCRmyPDF will make the PDF searchable
                cmd = ["ocrmypdf", "--skip-text", str(input_path), str(ocr_pdf)]
                self.run_subprocess(cmd, timeout=120)
                if ocr_pdf.exists():
                    pdf_source = ocr_pdf
                    logger.info("OCRmyPDF completed successfully. Swapping source to OCR'd PDF.")
            except Exception as e:
                logger.warning(f"OCRmyPDF failed: {e}. Falling back to standard processing.")

        try:
            # Route to the appropriate sub-pipeline
            if to_ext == ".docx":
                return self._to_docx(pdf_source, output_path, quality=quality)
            elif to_ext == ".pptx":
                return self._to_pptx(pdf_source, output_path, quality=quality)
            elif to_ext in [".xlsx", ".csv"]:
                return self._to_spreadsheet(pdf_source, output_path, to_ext)
            elif to_ext == ".txt":
                return self._to_txt(pdf_source, output_path)
            elif to_ext in image_targets:
                return self._to_image(pdf_source, output_path, to_ext, quality=quality)
            elif to_ext == ".html":
                return self._to_html(pdf_source, output_path)
            elif to_ext == ".zip":
                return self._to_zip(pdf_source, output_path, quality=quality)
            else:
                raise ValueError(f"PDFHandler does not support target: {to_ext}")
        finally:
            # Cleanup temp OCR PDF if created
            if is_scanned and pdf_source.name.endswith("_ocr.pdf") and pdf_source.exists():
                try:
                    pdf_source.unlink()
                except Exception as e:
                    logger.warning(f"Failed to delete temp OCR PDF: {e}")
            normalized_pdf = input_path.parent / f"{input_path.stem}_normalized.pdf"
            if normalized_pdf.exists() and normalized_pdf != input_path:
                try:
                    normalized_pdf.unlink()
                except Exception:
                    pass

    def _normalize_pdf(self, input_path: Path) -> Path:
        normalized = input_path.parent / f"{input_path.stem}_normalized.pdf"
        try:
            self.run_subprocess(
                ["qpdf", "--linearize", str(input_path), str(normalized)],
                timeout=60,
            )
            if normalized.exists() and normalized.stat().st_size > 0:
                return normalized
        except Exception as e:
            logger.warning(f"qpdf normalization skipped: {e}")
        return input_path

    def _is_scanned_pdf(self, pdf_path: Path) -> bool:
        """
        Determines if a PDF is scanned by checking the ratio of extractable text to pages.
        """
        try:
            with pdfplumber.open(pdf_path) as pdf:
                total_chars = 0
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    total_chars += len(text.strip())
                    if total_chars > 100 or i >= 4:
                        break
                return total_chars < 40
        except Exception as e:
            logger.error(f"Error checking if PDF is scanned: {e}")
            return False

    def _normalize_quality(self, quality: str) -> str:
        quality = (quality or "fast").strip().lower()
        return quality if quality in QUALITY_PRESETS else "fast"

    def _dpi(self, quality: str, key: str) -> int:
        return int(QUALITY_PRESETS.get(quality, QUALITY_PRESETS["fast"])[key])

    def _image_save_kwargs(self, to_ext: str, quality: str) -> dict:
        if to_ext in [".jpg", ".jpeg"]:
            return {"quality": 82} if quality == "fast" else {"quality": 95}
        if to_ext == ".webp":
            return {"quality": 78, "method": 4} if quality == "fast" else {"quality": 95, "method": 6}
        return {}

    def _to_docx(self, pdf_path: Path, output_path: Path, quality: str = "fast") -> Path:
        """
        Extracts semantic text paragraphs and compiles an editable Word Document.
        """
        # NOTE:
        # LibreOffice PDF import can generate VML-heavy DOCX that renders as blank/black
        # in some production viewers. Use deterministic text/image reconstruction only.
        logger.info("Using deterministic PDF->DOCX reconstruction pipeline (no LibreOffice PDF import).")

        logger.info("Extracting paragraphs using pdfplumber fallback...")
        doc = Document()
        has_text = False
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    has_text = True
                    for line in text.split("\n"):
                        line_strip = line.strip()
                        if line_strip:
                            doc.add_paragraph(line_strip)
        if has_text:
            doc.save(output_path)
            return output_path

        logger.warning("No extractable text found. Building image-based DOCX fallback to avoid blank output.")
        self._pdf_pages_to_docx_images(pdf_path, output_path, quality=quality)
        if self._docx_has_content(output_path):
            return output_path

        raise RuntimeError("Failed to convert PDF to DOCX using available pipelines.")

    def _docx_has_content(self, docx_path: Path) -> bool:
        """
        Quick quality gate: verifies that generated DOCX contains meaningful body text
        or embedded raster image media.
        """
        try:
            with zipfile.ZipFile(docx_path, "r") as zf:
                xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")
                names = [e.filename.lower() for e in zf.infolist()]
            text_nodes = re.findall(r"<w:t[^>]*>(.*?)</w:t>", xml, flags=re.DOTALL)
            meaningful_text = "".join(text_nodes).strip()
            has_raster_media = any(
                name.startswith("word/media/") and name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"))
                for name in names
            )
            # Drawing-only DOCX shells can still render as blank in some viewers.
            # Accept when we have meaningful text or concrete raster media.
            return len(meaningful_text) >= 20 or has_raster_media
        except Exception as e:
            logger.warning(f"DOCX validation check failed for {docx_path}: {e}")
            return False

    def _docx_has_problematic_vml_background(self, docx_path: Path) -> bool:
        """
        Detects a known LibreOffice PDF-import artifact where a full-page VML shape/fill
        can render as black pages in some Word viewers/editors.
        """
        try:
            with zipfile.ZipFile(docx_path, "r") as zf:
                xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")

            has_vml_shape = "<v:shape" in xml
            has_black_fill_hint = ('color2="black"' in xml) or ('<v:fill' in xml and 'black' in xml)
            # If this VML artifact appears and text content is sparse, prefer fallback output.
            text_nodes = re.findall(r"<w:t[^>]*>(.*?)</w:t>", xml, flags=re.DOTALL)
            text_len = len("".join(text_nodes).strip())
            return has_vml_shape and has_black_fill_hint and text_len < 200
        except Exception as e:
            logger.warning(f"Problematic VML detection failed for {docx_path}: {e}")
            return False

    def _pdf_pages_to_docx_images(self, pdf_path: Path, output_path: Path, quality: str = "fast") -> None:
        """
        Rasterizes PDF pages and embeds them into DOCX as a non-blank fallback when
        semantic text extraction is unavailable.
        """
        temp_img_dir = pdf_path.parent / "temp_docx_images"
        temp_img_dir.mkdir(parents=True, exist_ok=True)
        try:
            cmd = ["pdftoppm", "-png", "-r", str(self._dpi(quality, "docx_fallback_dpi")), str(pdf_path), str(temp_img_dir / "page")]
            self.run_subprocess(cmd, timeout=90)
            images = sorted(temp_img_dir.glob("page-*.png"))
            if not images:
                raise FileNotFoundError("No images generated for image-based DOCX fallback.")

            doc = Document()
            for idx, image_path in enumerate(images):
                if idx > 0:
                    doc.add_page_break()
                doc.add_picture(str(image_path), width=DocxInches(6.5))
            doc.save(output_path)
        finally:
            if temp_img_dir.exists():
                shutil.rmtree(temp_img_dir, ignore_errors=True)

    def _to_pptx(self, pdf_path: Path, output_path: Path, quality: str = "fast") -> Path:
        """
        Converts PDF pages to PNGs, then sets each page as a full-bleed slide in a Presentation.
        """
        logger.info("Converting PDF to images for slide deck construction...")
        temp_img_dir = pdf_path.parent / "temp_pptx_images"
        temp_img_dir.mkdir(parents=True, exist_ok=True)

        try:
            # 1. Extract PDF pages to PNGs
            cmd = ["pdftoppm", "-png", "-r", str(self._dpi(quality, "pptx_dpi")), str(pdf_path), str(temp_img_dir / "page")]
            self.run_subprocess(cmd, timeout=60)

            extracted_images = sorted(list(temp_img_dir.glob("page-*.png")))
            if not extracted_images:
                raise FileNotFoundError("pdftoppm failed to generate page images.")

            # 2. Build presentation using python-pptx
            prs = Presentation()
            # Match slide aspect ratio to first PDF page image to reduce stretching/cropping.
            with Image.open(extracted_images[0]) as first_img:
                width_px, height_px = first_img.size
            target_height = Inches(7.5)
            target_width = int(target_height * (width_px / max(height_px, 1)))
            prs.slide_width = target_width
            prs.slide_height = target_height
            blank_slide_layout = prs.slide_layouts[6] # Blank slide layout

            for img_path in extracted_images:
                slide = prs.slides.add_slide(blank_slide_layout)
                # Add full bleed background image
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(0), Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )

            prs.save(output_path)
            logger.info(f"Presentation saved successfully to {output_path}")
            return output_path

        finally:
            if temp_img_dir.exists():
                shutil.rmtree(temp_img_dir)

    def _to_spreadsheet(self, pdf_path: Path, output_path: Path, to_ext: str) -> Path:
        """
        Extracts tabular datasets using Camelot, with a robust fallback to pdfplumber.
        """
        import pandas as pd
        logger.info("Extracting tables using Camelot...")
        tables_df = []

        try:
            # Camelot lattice mode (works well for bordered grids)
            tables = camelot.read_pdf(str(pdf_path), pages="all", flavor="lattice", line_scale=40)
            if not tables or len(tables) == 0:
                # Retry stream mode for borderless tables
                tables = camelot.read_pdf(str(pdf_path), pages="all", flavor="stream", row_tol=10, column_tol=10)
            
            for t in tables:
                tables_df.append(t.df)
        except Exception as e:
            logger.warning(f"Camelot table extraction failed: {e}. Trying pdfplumber fallback...")

        # Fallback to pdfplumber table extractor
        if not tables_df:
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    for page in pdf.pages:
                        extracted = page.extract_tables()
                        for tbl in extracted:
                            if tbl:
                                # Clean None or empty headers
                                tables_df.append(pd.DataFrame(tbl))
            except Exception as e:
                logger.error(f"pdfplumber table extraction failed: {e}")

        # If no tables extracted, return a simple text-wrapped spreadsheet containing raw text lines
        if not tables_df:
            logger.warning("No structured tables detected. Creating a sheet with text paragraph lines.")
            paragraphs = []
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        paragraphs.extend([[line] for line in text.split("\n")])
            tables_df.append(pd.DataFrame(paragraphs, columns=["Extracted Content"]))

        # Normalize noisy extraction output so XLSX resembles polished converters:
        # - drop empty rows/cols
        # - remove synthetic numeric header row (0..n)
        # - collapse to compact key/value shape when appropriate
        normalized_tables = [self._normalize_extracted_table(df) for df in tables_df]

        # Save to XLSX or CSV
        if to_ext == ".xlsx":
            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                for idx, df in enumerate(normalized_tables):
                    sheet_name = f"Table_{idx + 1}"
                    # Excel sheet names capped at 31 chars
                    df.to_excel(sheet_name=sheet_name[:30], excel_writer=writer, index=False, header=False)

                # Post-format workbook for cleaner CloudConvert-like readability.
                self._postformat_xlsx(writer.book)
        else:
            # CSV - Merge all extracted tables and save
            merged = (
                pd.concat(normalized_tables, ignore_index=True)
                if len(normalized_tables) > 1
                else normalized_tables[0]
            )
            merged.to_csv(output_path, index=False)

        logger.info(f"Spreadsheet saved successfully to {output_path}")
        return output_path

    def _normalize_extracted_table(self, df):
        import pandas as pd

        cleaned = df.copy()

        def _clean_cell(v):
            if v is None:
                return ""
            s = str(v)
            s = s.replace("\r", " ").replace("\n", " ").replace("\t", " ")
            s = re.sub(r"\s+", " ", s).strip()
            return s

        cleaned = cleaned.map(_clean_cell)

        # Drop fully empty rows/cols.
        cleaned = cleaned.replace("", pd.NA)
        cleaned = cleaned.dropna(axis=0, how="all").dropna(axis=1, how="all")
        cleaned = cleaned.fillna("")

        if cleaned.empty:
            return df

        # Remove synthetic numeric header rows (e.g., 0,1,2,3,4) produced by dataframe defaults.
        first_row = [str(v).strip() for v in cleaned.iloc[0].tolist()]
        if first_row and all(v.isdigit() for v in first_row):
            numeric = [int(v) for v in first_row]
            if numeric == list(range(len(numeric))):
                cleaned = cleaned.iloc[1:].reset_index(drop=True)

        # If many sparse columns exist, compact to first 1-2 meaningful cells per row.
        if cleaned.shape[1] > 2:
            rows = []
            for _, row in cleaned.iterrows():
                vals = [str(v).strip() for v in row.tolist() if str(v).strip()]
                if not vals:
                    continue
                if len(vals) == 1:
                    rows.append([vals[0], ""])
                else:
                    rows.append([vals[0], vals[1]])
            if rows:
                cleaned = pd.DataFrame(rows, columns=["A", "B"])

        cleaned = cleaned.reset_index(drop=True)

        # PAN-document-aware canonical shape (editable cells, cloudconvert-like).
        canonical = self._build_pan_record_table(cleaned)
        if canonical is not None:
            return canonical

        # Generic normalization for non-PAN docs.
        if cleaned.shape[1] >= 2:
            key_col = cleaned.columns[0]
            val_col = cleaned.columns[1]
            cleaned[val_col] = [str(v).strip() for v in cleaned[val_col].tolist()]
        return cleaned

    def _build_pan_record_table(self, cleaned):
        import pandas as pd

        # Flatten rows into readable logical lines.
        lines = []
        for _, row in cleaned.iterrows():
            vals = [str(v).strip() for v in row.tolist() if str(v).strip()]
            if not vals:
                continue
            lines.append(" ".join(vals))

        if not lines:
            return None

        full_text = " ".join(lines).upper()
        if "PAN VERIFICATION RECORD" not in full_text:
            return None

        # Remove signature/noise lines that cloud-style XLSX omits from the core table.
        keep_lines = []
        for line in lines:
            u = line.upper()
            if u.startswith("DIGITALLY SIGNED"):
                continue
            if u.startswith("NOTE"):
                continue
            if "POWERED BY TCPDF" in u:
                continue
            if "(CID:" in u:
                continue
            keep_lines.append(line)

        pan = ""
        name = ""
        gender = ""
        dob = ""
        verified = ""
        has_title = False
        has_pan_label = False

        for line in keep_lines:
            u = line.upper()
            if "PAN VERIFICATION RECORD" in u:
                has_title = True
            if "PERMANENT ACCOUNT NUMBER" in u:
                has_pan_label = True

            m = re.search(r"\b([A-Z]{5}[0-9]{4}[A-Z])\b", u)
            if m:
                pan = m.group(1)

            if u.startswith("NAME "):
                name = line[5:].strip()
            elif u == "NAME":
                name = ""

            if u.startswith("GENDER "):
                gender = line[7:].strip()

            if u.startswith("DATE OF BIRTH "):
                dob = line[14:].strip()

            if u.startswith("VERIFIED ON "):
                verified = line[12:].strip()

        if not has_title or not has_pan_label:
            return None

        # Preserve user-friendly display format like cloud output.
        dob_dt = self._try_parse_excel_datetime(dob)
        if dob_dt:
            dob = dob_dt.strftime("%d-%m-%Y")
        ver_dt = self._try_parse_excel_datetime(verified)
        if ver_dt:
            verified = ver_dt.strftime("%d-%m-%Y %H:%M:%S")

        rows = [
            ["PAN VERIFICATION RECORD", ""],
            ["Permanent Account Number", ""],
            [pan, ""],
            ["NAME", name],
            ["GENDER", gender],
            ["DATE OF BIRTH", dob],
            ["VERIFIED ON", verified],
        ]
        return pd.DataFrame(rows, columns=["A", "B"])

    def _try_parse_excel_datetime(self, value: str):
        if not value:
            return None
        for fmt in ("%d-%m-%Y", "%d/%m/%Y", "%d/%m/%Y %H:%M:%S", "%d-%m-%Y %H:%M:%S"):
            try:
                return datetime.strptime(value, fmt)
            except Exception:
                continue
        return None

    def _postformat_xlsx(self, workbook):
        from openpyxl.styles import Alignment, Font, PatternFill, Border, Side

        for ws in workbook.worksheets:
            # Match cloud-like geometry.
            ws.column_dimensions["A"].width = 22
            ws.column_dimensions["B"].width = 60
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=2):
                for cell in row:
                    cell.alignment = Alignment(vertical="center", wrap_text=True)

            # Canonical PAN block styling.
            is_pan = str(ws.cell(1, 1).value or "").upper().strip() == "PAN VERIFICATION RECORD"
            if is_pan:
                # Merge title rows.
                for r in (1, 2, 3):
                    ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=2)

                ws.cell(1, 1).alignment = Alignment(horizontal="center", vertical="center")
                ws.cell(2, 1).alignment = Alignment(horizontal="center", vertical="center")
                ws.cell(3, 1).alignment = Alignment(horizontal="center", vertical="center")

                ws.cell(1, 1).font = Font(bold=False, size=26)
                ws.cell(2, 1).font = Font(bold=False, size=16)
                ws.cell(3, 1).font = Font(bold=False, size=28)

                for r in range(4, 8):
                    ws.cell(r, 1).font = Font(bold=True, size=16)
                    ws.cell(r, 2).font = Font(bold=False, size=16)

                # Light gray block look like cloud output.
                gray = PatternFill(fill_type="solid", fgColor="E8ECEC")
                for r in range(1, 8):
                    for c in (1, 2):
                        ws.cell(r, c).fill = gray

                # Thin grid border around A1:B7.
                thin = Side(style="thin", color="B7B7B7")
                border = Border(left=thin, right=thin, top=thin, bottom=thin)
                for r in range(1, 8):
                    for c in (1, 2):
                        ws.cell(r, c).border = border

                # Row heights for visual hierarchy.
                ws.row_dimensions[1].height = 36
                ws.row_dimensions[2].height = 28
                ws.row_dimensions[3].height = 40
                for r in range(4, 8):
                    ws.row_dimensions[r].height = 38
            else:
                # Generic fallback styling.
                for r in range(1, ws.max_row + 1):
                    a = ws.cell(r, 1).value
                    b = ws.cell(r, 2).value
                    if a and (b is None or str(b).strip() == "") and str(a).strip():
                        ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=2)
                        ws.cell(r, 1).font = Font(bold=True)

    def _to_txt(self, pdf_path: Path, output_path: Path) -> Path:
        """
        Extracts plain text.
        """
        content = []
        try:
            with pdfplumber.open(pdf_path) as pdf:
                for idx, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text:
                        content.append(text)
        except Exception as e:
            logger.warning(f"pdfplumber text extraction failed: {e}")

        if not content:
            try:
                txt = self.run_subprocess(
                    ["mutool", "draw", "-F", "txt", str(pdf_path)],
                    timeout=60,
                )
                if txt.strip():
                    content.append(txt.strip())
            except Exception as e:
                logger.warning(f"mutool text fallback failed: {e}")

        output_path.write_text("\n\n--- PAGE BREAK ---\n\n".join(content), encoding="utf-8")
        return output_path

    def _to_image(self, pdf_path: Path, output_path: Path, to_ext: str, quality: str = "fast") -> Path:
        """
        Converts PDF to PNG/JPEG/WEBP/HEIC.
        For multi-page PDFs, returns a ZIP containing one image per page.
        """
        temp_img_dir = pdf_path.parent / "temp_pdf_images"
        temp_img_dir.mkdir(parents=True, exist_ok=True)

        try:
            # 1. Render all pages to PNG
            cmd = ["pdftoppm", "-png", "-r", str(self._dpi(quality, "image_dpi")), str(pdf_path), str(temp_img_dir / "page")]
            self.run_subprocess(cmd, timeout=PDF_IMAGE_RENDER_TIMEOUT_SECONDS)

            extracted = sorted(list(temp_img_dir.glob("page-*.png")))
            if not extracted:
                raise FileNotFoundError("pdftoppm failed to generate page images.")

            # Multi-page PDF -> ZIP package so users receive all pages.
            if len(extracted) > 1:
                zip_output = output_path.with_suffix(".zip")
                with zipfile.ZipFile(zip_output, "w", zipfile.ZIP_DEFLATED) as zf:
                    for idx, png_path in enumerate(extracted, start=1):
                        page_output_name = f"page-{idx:03d}{to_ext if to_ext != '.jpeg' else '.jpg'}"
                        if to_ext == ".png":
                            zf.write(png_path, arcname=page_output_name)
                            continue

                        if to_ext == ".heic":
                            page_img_path = temp_img_dir / page_output_name
                            try:
                                cmd = ["heif-enc", str(png_path), str(page_img_path)]
                                self.run_subprocess(cmd, timeout=30)
                            except Exception:
                                if not HEIF_PLUGIN_ENABLED:
                                    raise
                                img = Image.open(png_path)
                                if img.mode in ("RGBA", "P"):
                                    img = img.convert("RGB")
                                img.save(page_img_path, format="HEIF", quality=90)
                            zf.write(page_img_path, arcname=page_output_name)
                            continue

                        # JPEG / WEBP conversion path
                        page_img_path = temp_img_dir / page_output_name
                        img = Image.open(png_path)
                        if to_ext in [".jpg", ".jpeg"] and img.mode in ("RGBA", "P"):
                            img = img.convert("RGB")
                        img.save(page_img_path, **self._image_save_kwargs(to_ext, quality))
                        zf.write(page_img_path, arcname=page_output_name)
                return zip_output

            png_page = extracted[0]

            if to_ext == ".png":
                shutil.move(str(png_page), str(output_path))
            elif to_ext == ".heic":
                # Convert PNG to HEIC using heif-enc
                try:
                    cmd = ["heif-enc", str(png_page), str(output_path)]
                    self.run_subprocess(cmd, timeout=30)
                except Exception:
                    if not HEIF_PLUGIN_ENABLED:
                        raise
                    img = Image.open(png_page)
                    if img.mode in ("RGBA", "P"):
                        img = img.convert("RGB")
                    img.save(output_path, format="HEIF", quality=90)
            else:
                # Pillow convert for JPG, JPEG, WEBP
                img = Image.open(png_page)
                if to_ext in [".jpg", ".jpeg"] and img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                img.save(output_path, **self._image_save_kwargs(to_ext, quality))

            return output_path
        finally:
            if temp_img_dir.exists():
                shutil.rmtree(temp_img_dir)

    def _to_html(self, pdf_path: Path, output_path: Path) -> Path:
        """
        Generates high-fidelity HTML, preferring layout-preserving engines first:
        1) pdf2htmlEX
        2) poppler pdftohtml (single-file, CSS-driven)
        3) text fallback template
        """
        # 1) Best fidelity: pdf2htmlEX (very close to CloudConvert style output)
        if self._try_pdf2htmlex(pdf_path, output_path):
            return output_path

        # 2) Good fallback: poppler's pdftohtml with embedded CSS and no frame split
        if self._try_pdftohtml(pdf_path, output_path):
            return output_path

        # 3) Fidelity fallback: render each PDF page as image and wrap in HTML.
        # This preserves layout/branding even when html converters are unavailable.
        if self._to_html_as_images(pdf_path, output_path):
            return output_path

        # 4) Last fallback: semantic text extraction template
        html_blocks = []
        with pdfplumber.open(pdf_path) as pdf:
            for idx, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    paragraphs = page_text.split("\n")
                    formatted = "".join(
                        f"<p style='margin-bottom:12px;'>{html.escape(p)}</p>"
                        for p in paragraphs
                        if p.strip()
                    )
                    html_blocks.append(f"<div class='page' style='margin: 40px auto; padding: 20px; border: 1px solid #ccc; max-width: 800px; background:#fff;'><h2>Page {idx + 1}</h2>{formatted}</div>")

        body_content = "\n".join(html_blocks)
        html_template = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>Converted Document - {pdf_path.stem}</title>
    <style>
        body {{ font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif; background-color: #f7f9fc; color: #333; margin: 0; padding: 20px; }}
        h2 {{ color: #4F46E5; border-bottom: 2px solid #E5E7EB; padding-bottom: 8px; }}
    </style>
</head>
<body>
    {body_content}
</body>
</html>
"""
        output_path.write_text(html_template, encoding="utf-8")
        return output_path

    def _to_html_as_images(self, pdf_path: Path, output_path: Path) -> bool:
        temp_img_dir = pdf_path.parent / "temp_html_images"
        temp_img_dir.mkdir(parents=True, exist_ok=True)
        try:
            # 170 DPI keeps text crisp without huge HTML output.
            cmd = ["pdftoppm", "-png", "-r", "170", str(pdf_path), str(temp_img_dir / "page")]
            self.run_subprocess(cmd, timeout=PDF_IMAGE_RENDER_TIMEOUT_SECONDS)

            pages = sorted(temp_img_dir.glob("page-*.png"))
            if not pages:
                return False

            page_nodes = []
            for p in pages:
                b64 = base64.b64encode(p.read_bytes()).decode("ascii")
                page_nodes.append(
                    "<div class='page'><img alt='PDF page' "
                    f"src='data:image/png;base64,{b64}'/></div>"
                )

            html_out = """<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <title>Converted Document</title>
  <style>
    body {
      margin: 0;
      padding: 22px 0;
      background: #d8dbe0;
      font-family: Arial, sans-serif;
    }
    .page {
      width: fit-content;
      margin: 0 auto 18px auto;
      background: #fff;
      box-shadow: 0 1px 3px rgba(0, 0, 0, 0.35);
    }
    .page img {
      display: block;
      max-width: min(96vw, 1100px);
      height: auto;
    }
  </style>
</head>
<body>
""" + "\n".join(page_nodes) + """
</body>
</html>"""
            output_path.write_text(html_out, encoding="utf-8")
            return True
        except Exception as e:
            logger.info(f"Image-based HTML fallback failed: {e}")
            return False
        finally:
            if temp_img_dir.exists():
                shutil.rmtree(temp_img_dir, ignore_errors=True)

    def _try_pdf2htmlex(self, pdf_path: Path, output_path: Path) -> bool:
        try:
            self.run_subprocess(
                [
                    "pdf2htmlEX",
                    "--embed-css",
                    "1",
                    "--embed-font",
                    "1",
                    "--embed-image",
                    "1",
                    "--embed-javascript",
                    "1",
                    "--embed-outline",
                    "0",
                    "--dest-dir",
                    str(output_path.parent),
                    str(pdf_path),
                    output_path.name,
                ],
                timeout=180,
            )
            if not (output_path.exists() and output_path.stat().st_size > 0):
                return False
            self._inline_local_assets(output_path)
            return not self._looks_broken_html(output_path)
        except Exception as e:
            logger.info(f"pdf2htmlEX unavailable/failed, trying pdftohtml fallback: {e}")
            return False

    def _try_pdftohtml(self, pdf_path: Path, output_path: Path) -> bool:
        try:
            # pdftohtml writes to <prefix>.html, so pass path without suffix as prefix.
            prefix = output_path.with_suffix("")
            self.run_subprocess(
                [
                    "pdftohtml",
                    "-s",          # single HTML file
                    "-c",          # complex output (positioned text)
                    "-hidden",     # include hidden text (better fidelity/searchability)
                    "-noframes",   # no frame split
                    "-enc",
                    "UTF-8",
                    str(pdf_path),
                    str(prefix),
                ],
                timeout=120,
            )
            produced = prefix.with_suffix(".html")
            if produced.exists() and produced.stat().st_size > 0:
                if produced != output_path:
                    shutil.move(str(produced), str(output_path))
                self._inline_local_assets(output_path)
                return not self._looks_broken_html(output_path)
            return False
        except Exception as e:
            logger.info(f"pdftohtml unavailable/failed, falling back to text-template HTML: {e}")
            return False

    def _looks_broken_html(self, html_path: Path) -> bool:
        """
        Heuristic: detect common broken-render output where text layers exist but
        referenced background/image assets are missing.
        """
        try:
            src = html_path.read_text(encoding="utf-8", errors="ignore").lower()
            if "background image" in src:
                return True
            # Has image tags but none are embedded/absolute after inlining attempt.
            img_refs = re.findall(r"""<img[^>]+src=["']([^"']+)["']""", src, flags=re.I)
            if img_refs and all((not s.startswith("data:")) for s in img_refs):
                return True
            return False
        except Exception:
            return True

    def _inline_local_assets(self, html_path: Path) -> None:
        """
        Replace local img/src and stylesheet href assets with data-URIs so the HTML
        remains portable and does not break when moved between folders/machines.
        """
        try:
            src = html_path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return

        base_dir = html_path.parent

        def to_data_uri(asset_rel: str) -> str:
            if asset_rel.startswith(("data:", "http://", "https://", "//", "#")):
                return asset_rel
            asset_path = (base_dir / asset_rel).resolve()
            if not asset_path.exists() or not asset_path.is_file():
                return asset_rel
            mime, _ = mimetypes.guess_type(str(asset_path))
            mime = mime or "application/octet-stream"
            raw = asset_path.read_bytes()
            b64 = base64.b64encode(raw).decode("ascii")
            return f"data:{mime};base64,{b64}"

        def replace_img(m):
            before, url, after = m.group(1), m.group(2), m.group(3)
            return f'{before}{to_data_uri(url)}{after}'

        def replace_link_css(m):
            whole, href = m.group(0), m.group(1)
            data = to_data_uri(href)
            if data == href:
                return whole
            # Convert CSS file to an inline <style> when possible.
            if data.startswith("data:text/css;base64,"):
                try:
                    css_raw = base64.b64decode(data.split(",", 1)[1]).decode("utf-8", errors="ignore")
                    return f"<style>{css_raw}</style>"
                except Exception:
                    return whole
            return whole

        src = re.sub(
            r"""(<img[^>]+src=["'])([^"']+)(["'])""",
            replace_img,
            src,
            flags=re.I,
        )
        src = re.sub(
            r"""<link[^>]+rel=["']stylesheet["'][^>]+href=["']([^"']+)["'][^>]*>""",
            replace_link_css,
            src,
            flags=re.I,
        )

        try:
            html_path.write_text(src, encoding="utf-8")
        except Exception:
            pass

    def _to_zip(self, pdf_path: Path, output_path: Path, quality: str = "fast") -> Path:
        """
        Converts all pages of a PDF to high-fidelity images, packaging them neatly into a single ZIP.
        """
        temp_img_dir = pdf_path.parent / "temp_zip_images"
        temp_img_dir.mkdir(parents=True, exist_ok=True)

        try:
            # 1. Render all pages
            cmd = ["pdftoppm", "-png", "-r", str(self._dpi(quality, "zip_dpi")), str(pdf_path), str(temp_img_dir / "page")]
            self.run_subprocess(cmd, timeout=60)

            extracted = list(temp_img_dir.glob("page-*.png"))
            if not extracted:
                raise FileNotFoundError("pdftoppm failed to extract pages.")

            # 2. Package to ZIP
            with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zip_file:
                for img_path in sorted(extracted):
                    zip_file.write(img_path, arcname=img_path.name)

            return output_path
        finally:
            if temp_img_dir.exists():
                shutil.rmtree(temp_img_dir)
